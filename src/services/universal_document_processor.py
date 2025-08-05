"""
Universal Document Processing Service
Handles any document type: PDF, Word, Excel, PowerPoint, Text, HTML, etc.
"""

import os
import logging
import mimetypes
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import hashlib
from datetime import datetime

# Document processing imports
import PyPDF2
try:
    import pypdf
except ImportError:
    pypdf = None

from docx import Document as DocxDocument
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

from bs4 import BeautifulSoup

from src.models.schemas import DocumentChunk
from src.core.config import settings

logger = logging.getLogger(__name__)


class UniversalDocumentProcessor:
    """Process any document type and extract text content"""
    
    def __init__(self):
        self.supported_types = {
            'application/pdf': self._process_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_xlsx,
            'text/plain': self._process_text,
            'text/html': self._process_html,
            'text/markdown': self._process_text,
            'application/json': self._process_json,
            'text/csv': self._process_csv,
        }
    
    def process_file(self, file_path: str, file_content: bytes = None) -> List[DocumentChunk]:
        """
        Process any file and extract text chunks
        
        Args:
            file_path: Path to the file or filename
            file_content: Raw file content (for uploaded files)
        
        Returns:
            List of DocumentChunk objects
        """
        try:
            # Determine file type
            mime_type = self._detect_file_type(file_path, file_content)
            logger.info(f"Processing file: {file_path}, type: {mime_type}")
            
            # Get processor for this file type
            processor = self.supported_types.get(mime_type)
            if not processor:
                # Try to process as text if unknown type
                logger.warning(f"Unknown file type {mime_type}, trying as text")
                processor = self._process_text
            
            # Extract text content
            if file_content:
                text_content = processor(file_content, file_path)
            else:
                with open(file_path, 'rb') as f:
                    text_content = processor(f.read(), file_path)
            
            # Create document chunks
            chunks = self._create_chunks(text_content, file_path, mime_type)
            
            logger.info(f"✅ Processed {file_path}: {len(chunks)} chunks extracted")
            return chunks
            
        except Exception as e:
            logger.error(f"❌ Failed to process {file_path}: {str(e)}")
            return []
    
    def _detect_file_type(self, file_path: str, file_content: bytes = None) -> str:
        """Detect file MIME type"""
        # First try by extension
        mime_type, _ = mimetypes.guess_type(file_path)
        
        if mime_type:
            return mime_type
        
        # Try by file content if available
        if file_content:
            # Check for common file signatures
            if file_content.startswith(b'%PDF'):
                return 'application/pdf'
            elif file_content.startswith(b'PK\x03\x04'):
                # ZIP-based formats (docx, xlsx, pptx)
                if '.docx' in file_path.lower():
                    return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                elif '.xlsx' in file_path.lower():
                    return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                elif '.pptx' in file_path.lower():
                    return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        # Default to text
        return 'text/plain'
    
    def _process_pdf(self, content: bytes, filename: str) -> str:
        """Extract text from PDF"""
        try:
            from io import BytesIO
            
            # Try pypdf first (better for newer PDFs)
            if pypdf:
                try:
                    pdf_reader = pypdf.PdfReader(BytesIO(content))
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    
                    if text.strip():
                        return text
                except Exception:
                    pass
            
            # Fallback to PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text
            
        except Exception as e:
            logger.warning(f"PDF processing failed: {e}")
            return f"PDF document: {filename} (text extraction failed)"
    
    def _process_docx(self, content: bytes, filename: str) -> str:
        """Extract text from Word document"""
        try:
            from io import BytesIO
            doc = DocxDocument(BytesIO(content))
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text
            
        except Exception as e:
            logger.warning(f"DOCX processing failed: {e}")
            return f"Word document: {filename} (text extraction failed)"
    
    def _process_pptx(self, content: bytes, filename: str) -> str:
        """Extract text from PowerPoint presentation"""
        if not Presentation:
            return f"PowerPoint presentation: {filename} (python-pptx not installed)"
        
        try:
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text
            
        except Exception as e:
            logger.warning(f"PPTX processing failed: {e}")
            return f"PowerPoint presentation: {filename} (text extraction failed)"
    
    def _process_xlsx(self, content: bytes, filename: str) -> str:
        """Extract text from Excel spreadsheet"""
        if not openpyxl:
            return f"Excel spreadsheet: {filename} (openpyxl not installed)"
        
        try:
            from io import BytesIO
            workbook = openpyxl.load_workbook(BytesIO(content))
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            
            return text
            
        except Exception as e:
            logger.warning(f"XLSX processing failed: {e}")
            return f"Excel spreadsheet: {filename} (text extraction failed)"
    
    def _process_text(self, content: bytes, filename: str) -> str:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # If all fail, use utf-8 with error handling
            return content.decode('utf-8', errors='replace')
            
        except Exception as e:
            logger.warning(f"Text processing failed: {e}")
            return f"Text document: {filename} (text extraction failed)"
    
    def _process_html(self, content: bytes, filename: str) -> str:
        """Extract text from HTML"""
        try:
            html_content = content.decode('utf-8', errors='replace')
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            logger.warning(f"HTML processing failed: {e}")
            return f"HTML document: {filename} (text extraction failed)"
    
    def _process_json(self, content: bytes, filename: str) -> str:
        """Extract text from JSON"""
        try:
            import json
            json_content = content.decode('utf-8', errors='replace')
            data = json.loads(json_content)
            
            # Convert JSON to readable text
            return json.dumps(data, indent=2, ensure_ascii=False)
            
        except Exception as e:
            logger.warning(f"JSON processing failed: {e}")
            return f"JSON document: {filename} (text extraction failed)"
    
    def _process_csv(self, content: bytes, filename: str) -> str:
        """Extract text from CSV"""
        try:
            import csv
            from io import StringIO
            
            csv_content = content.decode('utf-8', errors='replace')
            csv_reader = csv.reader(StringIO(csv_content))
            
            text = ""
            for row in csv_reader:
                text += " | ".join(row) + "\n"
            
            return text
            
        except Exception as e:
            logger.warning(f"CSV processing failed: {e}")
            return f"CSV document: {filename} (text extraction failed)"

    def _create_chunks(self, text: str, filename: str, mime_type: str) -> List[DocumentChunk]:
        """Split text into chunks and create DocumentChunk objects"""
        if not text or not text.strip():
            return []

        # Generate document ID
        doc_id = hashlib.md5(f"{filename}_{datetime.now().isoformat()}".encode()).hexdigest()

        # Split text into chunks
        chunks = []
        chunk_size = settings.MAX_CHUNK_SIZE
        overlap = settings.CHUNK_OVERLAP

        # Split by paragraphs first, then by sentences if needed
        paragraphs = text.split('\n\n')
        current_chunk = ""
        chunk_index = 0

        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue

            # If adding this paragraph would exceed chunk size
            if len(current_chunk) + len(paragraph) > chunk_size and current_chunk:
                # Save current chunk
                chunk = DocumentChunk(
                    chunk_id=f"{doc_id}_chunk_{chunk_index}",
                    document_id=doc_id,
                    content=current_chunk.strip(),
                    metadata={
                        'filename': filename,
                        'mime_type': mime_type,
                        'chunk_index': chunk_index,
                        'source': 'uploaded_document',
                        'processed_at': datetime.now().isoformat()
                    }
                )
                chunks.append(chunk)

                # Start new chunk with overlap
                if overlap > 0 and len(current_chunk) > overlap:
                    current_chunk = current_chunk[-overlap:] + "\n" + paragraph
                else:
                    current_chunk = paragraph
                chunk_index += 1
            else:
                # Add paragraph to current chunk
                if current_chunk:
                    current_chunk += "\n\n" + paragraph
                else:
                    current_chunk = paragraph

        # Add final chunk if there's content
        if current_chunk.strip():
            chunk = DocumentChunk(
                chunk_id=f"{doc_id}_chunk_{chunk_index}",
                document_id=doc_id,
                content=current_chunk.strip(),
                metadata={
                    'filename': filename,
                    'mime_type': mime_type,
                    'chunk_index': chunk_index,
                    'source': 'uploaded_document',
                    'processed_at': datetime.now().isoformat()
                }
            )
            chunks.append(chunk)

        return chunks

    def get_supported_types(self) -> List[str]:
        """Get list of supported file types"""
        return list(self.supported_types.keys())

    def is_supported(self, file_path: str) -> bool:
        """Check if file type is supported"""
        mime_type = self._detect_file_type(file_path)
        return mime_type in self.supported_types
